"""Build iLaunchify seed pitch deck (12 slides, 16:9 widescreen).

Brand:
  Pink     #FF2E63 brand
  Black    #0A0A0A primary surface
  Neon     #B5FF3D accent (dark surfaces only)
  Cream    #F3EFE8 light accent surface
  White    #FFFFFF
  Ink-700  #2B2B2B body on light
  Ink-500  #6B6B6B muted

Typography: header = Bricolage Grotesque (falls back to system), body = Inter.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# -------- palette --------
PINK = RGBColor(0xFF, 0x2E, 0x63)
BLACK = RGBColor(0x0A, 0x0A, 0x0A)
NEON = RGBColor(0xB5, 0xFF, 0x3D)
CREAM = RGBColor(0xF3, 0xEF, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK_900 = RGBColor(0x1A, 0x1A, 0x1A)
INK_700 = RGBColor(0x2B, 0x2B, 0x2B)
INK_500 = RGBColor(0x6B, 0x6B, 0x6B)
INK_300 = RGBColor(0xC9, 0xC4, 0xBC)
DARK_PANEL = RGBColor(0x16, 0x16, 0x16)

HEADER_FONT = "Bricolage Grotesque"
BODY_FONT = "Inter"

# -------- presentation setup --------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# -------- helpers --------

def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    *,
    font=BODY_FONT,
    size=14,
    color=INK_700,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    line_spacing=1.15,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    # support multiline by splitting on \n
    lines = text.split("\n") if isinstance(text, str) else [text]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = ln
        f = r.font
        f.name = font
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
    return tb


def add_rich(slide, runs, x, y, w, h, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    """runs: list of (text, dict_opts) where dict_opts can include font, size, color, bold, italic, newline."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    para = tf.paragraphs[0]
    para.alignment = align
    para.line_spacing = line_spacing
    first = True
    for text, opts in runs:
        if opts.get("newline_before") and not first:
            para = tf.add_paragraph()
            para.alignment = align
            para.line_spacing = line_spacing
        r = para.add_run()
        r.text = text
        f = r.font
        f.name = opts.get("font", BODY_FONT)
        f.size = Pt(opts.get("size", 14))
        f.bold = opts.get("bold", False)
        f.italic = opts.get("italic", False)
        f.color.rgb = opts.get("color", INK_700)
        first = False
    return tb


def add_rect(slide, x, y, w, h, fill, line_color=None, line_width=None, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line_color is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line_color
        if line_width is not None:
            s.line.width = line_width
    return s


def add_pill(slide, x, y, w, h, text, *, fill, text_color, size=11, bold=True, font=BODY_FONT):
    p = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    p.shadow.inherit = False
    p.adjustments[0] = 0.5
    p.fill.solid()
    p.fill.fore_color.rgb = fill
    p.line.fill.background()
    tf = p.text_frame
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    return p


def add_circle_number(slide, x, y, d, n, *, fill, text_color):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    c.shadow.inherit = False
    c.fill.solid()
    c.fill.fore_color.rgb = fill
    c.line.fill.background()
    tf = c.text_frame
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = str(n)
    r.font.name = HEADER_FONT
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = text_color
    return c


def set_speaker_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def add_footer(slide, *, light=True):
    # No footer chrome — caused overlaps with bottom-of-slide content.
    # Keep function as a no-op so existing call sites remain valid.
    return None


def add_brand_mark(slide, x, y, *, on_dark=False):
    """small wordmark."""
    text_color = WHITE if on_dark else BLACK
    accent = NEON if on_dark else PINK
    # dot
    dot = add_rect(slide, x, y + Inches(0.18), Inches(0.18), Inches(0.18), accent, shape=MSO_SHAPE.OVAL)
    add_text(
        slide,
        "iLaunchify",
        x + Inches(0.28),
        y,
        Inches(2.5),
        Inches(0.5),
        font=HEADER_FONT,
        size=20,
        bold=True,
        color=text_color,
        anchor=MSO_ANCHOR.TOP,
    )


# =========================================================
# SLIDE 1 — TITLE (light)
# =========================================================
s1 = prs.slides.add_slide(BLANK)
add_bg(s1, WHITE)

# Big pink corner mark
add_rect(s1, Inches(0), Inches(0), Inches(0.45), SH, PINK)

# Brandmark
add_brand_mark(s1, Inches(1.0), Inches(0.7))

# Cohort year pill
add_pill(s1, Inches(10.6), Inches(0.85), Inches(1.9), Inches(0.42), "SEED ROUND · 2026", fill=BLACK, text_color=WHITE, size=10)

# Headline
add_text(
    s1,
    "Production orchestration\nfor the next generation\nof CPG creators.",
    Inches(1.0),
    Inches(2.05),
    Inches(11.3),
    Inches(3.0),
    font=HEADER_FONT,
    size=54,
    bold=True,
    color=BLACK,
    line_spacing=1.05,
)

# Subhead — with italic Fraunces-style emphasis
add_rich(
    s1,
    [
        ("iLaunchify lets creators ship their own CPG brand ", {"size": 18, "color": INK_700}),
        ("without ", {"size": 18, "color": INK_700, "italic": True, "font": "Fraunces"}),
        ("owning a factory.", {"size": 18, "color": INK_700}),
    ],
    Inches(1.0),
    Inches(5.0),
    Inches(11.3),
    Inches(0.8),
)

# Founder block
add_text(s1, "Pavel Georgiev · Founder & CEO", Inches(1.0), Inches(6.35), Inches(7), Inches(0.3), size=13, bold=True, color=BLACK)
add_text(s1, "georgiev.pavel@gmail.com · June 2026", Inches(1.0), Inches(6.7), Inches(7), Inches(0.3), size=11, color=INK_500)

# Right-side neon decoration
add_rect(s1, Inches(11.4), Inches(6.35), Inches(1.5), Inches(0.65), NEON, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s1, "v1 built.\nbeta next.", Inches(11.4), Inches(6.4), Inches(1.5), Inches(0.6), size=10, bold=True, color=BLACK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=BODY_FONT)

set_speaker_notes(
    s1,
    "iLaunchify is the production orchestration layer for creator-led CPG. We let influencers and small DTC brand owners launch real supplement, "
    "functional food, or pet products without finding their own manufacturer, printer, co-packer, or warehouse — we decompose every order into a "
    "multi-partner workflow graph and hide the complexity. The platform is built — four Next.js apps, around 590 shipped tasks, FDA-aware compliance "
    "scanning, Stripe Connect end-to-end — and we are raising a seed round to run a structured beta and reach GA. This deck walks through the problem, "
    "the architectural moat we locked early, monetization, the beta plan, and the ask."
)

# =========================================================
# SLIDE 2 — THE PROBLEM (light)
# =========================================================
s2 = prs.slides.add_slide(BLANK)
add_bg(s2, WHITE)
add_brand_mark(s2, Inches(0.7), Inches(0.5))
add_pill(s2, Inches(0.7), Inches(1.2), Inches(1.3), Inches(0.36), "01 · PROBLEM", fill=BLACK, text_color=WHITE, size=10)

add_text(
    s2,
    "The CPG creator wave is real.\nThe production stack is broken.",
    Inches(0.7),
    Inches(1.75),
    Inches(12),
    Inches(1.7),
    font=HEADER_FONT,
    size=36,
    bold=True,
    color=BLACK,
    line_spacing=1.05,
)

# Left column: the wave
add_text(s2, "The headline names — and the long tail", Inches(0.7), Inches(3.55), Inches(5.6), Inches(0.4), size=12, bold=True, color=PINK, font=HEADER_FONT)
add_text(
    s2,
    "Mr Beast Feastables. Logan Paul Prime.\nEmma Chamberlain Coffee. Kim Kardashian Skims.",
    Inches(0.7),
    Inches(3.95),
    Inches(5.8),
    Inches(1.1),
    size=14,
    color=INK_700,
    line_spacing=1.35,
)
add_text(
    s2,
    "Below them sits a much larger cohort:\ncreators with 50K–2M followers, real channel\nrevenue, and a brand-able audience — and no\nway to actually ship a CPG product.",
    Inches(0.7),
    Inches(5.05),
    Inches(5.8),
    Inches(1.8),
    size=13,
    color=INK_700,
    line_spacing=1.35,
)

# Right column: pain checklist on cream panel
add_rect(s2, Inches(6.95), Inches(3.4), Inches(5.7), Inches(3.5), CREAM, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s2, "Today, launching means…", Inches(7.2), Inches(3.55), Inches(5.2), Inches(0.4), size=12, bold=True, color=BLACK, font=HEADER_FONT)

pain_items = [
    "Hunt and validate a contract manufacturer",
    "Source a label printer separately",
    "Find a co-packer that handles your format",
    "Negotiate a warehouse + fulfillment partner",
    "Manage FDA labeling compliance yourself",
    "Hand-stitch a checkout that ties it all together",
]
for i, item in enumerate(pain_items):
    row_y = Inches(3.95 + i * 0.42)
    add_rect(s2, Inches(7.2), row_y + Inches(0.13), Inches(0.12), Inches(0.12), PINK, shape=MSO_SHAPE.OVAL)
    add_text(s2, item, Inches(7.45), row_y, Inches(5), Inches(0.35), size=12, color=INK_700)

# bottom row stat
add_rect(s2, Inches(0.7), Inches(6.55), Inches(5.8), Inches(0.55), BLACK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_rich(
    s2,
    [
        ("12–18 months ", {"size": 14, "bold": True, "color": NEON, "font": HEADER_FONT}),
        ("to launch — or the creator gives up.", {"size": 13, "color": WHITE}),
    ],
    Inches(0.95), Inches(6.62), Inches(5.5), Inches(0.4),
)

add_footer(s2)

set_speaker_notes(
    s2,
    "The creator-led CPG wave is real and accelerating. Mr Beast launched Feastables, Logan Paul launched Prime, Emma Chamberlain launched coffee, "
    "Kim K launched Skims. Those are the headlines. The much larger opportunity is the cohort below them — 50K to 2M follower creators with real channel "
    "revenue and a brand-able audience, who today still can't actually ship a CPG product without a 12 to 18 month odyssey of finding and validating a "
    "contract manufacturer, sourcing a label printer separately, lining up a co-packer that handles their format, negotiating a warehouse partner, managing "
    "FDA labeling on their own, and stitching together a checkout. Most of them give up. The few who succeed are the ones with operator co-founders or "
    "agency money behind them. That gap is the wedge."
)

# =========================================================
# SLIDE 3 — THE INSIGHT (light)
# =========================================================
s3 = prs.slides.add_slide(BLANK)
add_bg(s3, WHITE)
add_brand_mark(s3, Inches(0.7), Inches(0.5))
add_pill(s3, Inches(0.7), Inches(1.2), Inches(1.3), Inches(0.36), "02 · INSIGHT", fill=BLACK, text_color=WHITE, size=10)

add_text(
    s3,
    "Three forces just crossed\nthe line at the same time.",
    Inches(0.7), Inches(1.75), Inches(12), Inches(1.7),
    font=HEADER_FONT, size=36, bold=True, color=BLACK, line_spacing=1.05,
)

# 3 numbered cards
card_y = Inches(3.55)
card_h = Inches(3.2)
card_w = Inches(3.95)
gap = Inches(0.2)
xs = [Inches(0.7), Inches(0.7) + card_w + gap, Inches(0.7) + (card_w + gap) * 2]

cards = [
    ("01", "Creators have monetization fluency.",
     "TikTok Shop, Shopify, and audience-economy capital didn't exist at this scale 24 months ago. The next layer of creators now has both the cash and the operating sophistication to launch a brand."),
    ("02", "CPG capacity has consolidated.",
     "Contract manufacturers exited 2024 with structurally hungry small-batch capacity. They want branded, smaller-run business — they just have no clean way to find creators who match."),
    ("03", "Design + compliance got accessible.",
     "Modern in-browser canvas tooling plus AI-aware FDA scanning crossed the line where a non-expert can produce print-ready, regulated CPG packaging. We could not have built this in 2022."),
]
for i, (n, title, body) in enumerate(cards):
    add_rect(s3, xs[i], card_y, card_w, card_h, WHITE, line_color=INK_300, line_width=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    # number badge
    add_rect(s3, xs[i] + Inches(0.3), card_y + Inches(0.3), Inches(0.6), Inches(0.6), PINK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s3, n, xs[i] + Inches(0.3), card_y + Inches(0.3), Inches(0.6), Inches(0.6), size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=HEADER_FONT)
    add_text(s3, title, xs[i] + Inches(0.3), card_y + Inches(1.05), card_w - Inches(0.6), Inches(0.85), size=17, bold=True, color=BLACK, font=HEADER_FONT, line_spacing=1.15)
    add_text(s3, body, xs[i] + Inches(0.3), card_y + Inches(1.95), card_w - Inches(0.6), Inches(1.15), size=11, color=INK_700, line_spacing=1.3)

# bottom takeaway
add_rich(
    s3,
    [
        ("None of these are sufficient alone. ", {"size": 14, "color": INK_500, "italic": True, "font": "Fraunces"}),
        ("All three together support a platform.", {"size": 14, "color": BLACK, "bold": True}),
    ],
    Inches(0.7), Inches(6.95), Inches(12), Inches(0.4),
)

set_speaker_notes(
    s3,
    "We didn't build this because the problem is new — it's been around since the first influencer launched a protein powder. We built it now because three "
    "structural conditions changed simultaneously in the last 24 months. First, audience-economy creators finally have the monetization fluency and capital — "
    "TikTok Shop, Shopify, paid memberships — to be real CPG buyers. Second, the post-2024 contract manufacturing landscape has consolidated capacity that's "
    "hungry for branded small-batch work but has no clean acquisition channel for it. Third, modern Fabric.js-grade browser canvas plus AI-tuned compliance "
    "scanning crossed the line where a non-expert can actually produce print-ready FDA-compliant CPG packaging. Any one of these isn't enough. The three together "
    "support a platform — and we're the first one designed for the intersection."
)

# =========================================================
# SLIDE 4 — THE PRODUCT (DARK HERO)
# =========================================================
s4 = prs.slides.add_slide(BLANK)
add_bg(s4, BLACK)
add_brand_mark(s4, Inches(0.7), Inches(0.5), on_dark=True)
add_pill(s4, Inches(0.7), Inches(1.2), Inches(1.5), Inches(0.36), "03 · PRODUCT", fill=NEON, text_color=BLACK, size=10)

add_rich(
    s4,
    [
        ("iLaunchify is the\n", {"size": 40, "color": WHITE, "bold": True, "font": HEADER_FONT}),
        ("production orchestration ", {"size": 40, "color": NEON, "bold": True, "font": HEADER_FONT}),
        ("layer.", {"size": 40, "color": WHITE, "bold": True, "font": HEADER_FONT}),
    ],
    Inches(0.7), Inches(1.75), Inches(12.2), Inches(2),
    line_spacing=1.05,
)

# 4 product pillars in a 2x2 grid
pillars = [
    ("Design Studio",
     "Fabric.js packaging canvas. Live FDA-aware compliance scan. Brand asset library that feeds directly into print-ready PDF export."),
    ("Marketplace",
     "13 product categories. 121 subcategories. Manufacturer-built templates filtered by format, region, capability, and certification."),
    ("Orchestration",
     "Each order decomposes into a multi-partner workflow graph: manufacturer → printer → co-packer → warehouse → creator's channel."),
    ("Money + manifest",
     "Stripe Connect multi-party payouts. Payment held until every partner approves. Manifest versioned and audit-logged end-to-end."),
]
gx = Inches(0.7); gy = Inches(3.65); gw = Inches(5.95); gh = Inches(1.4); gh_gap = Inches(0.15)
positions = [(0, 0), (1, 0), (0, 1), (1, 1)]  # col, row
for i, (title, body) in enumerate(pillars):
    col, row = positions[i]
    x = gx + (gw + Inches(0.3)) * col
    y = gy + (gh + gh_gap) * row
    add_rect(s4, x, y, gw, gh, DARK_PANEL, line_color=RGBColor(0x2A, 0x2A, 0x2A), line_width=Pt(0.5), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s4, title, x + Inches(0.3), y + Inches(0.18), gw - Inches(0.6), Inches(0.4), size=16, bold=True, color=NEON, font=HEADER_FONT)
    add_text(s4, body, x + Inches(0.3), y + Inches(0.6), gw - Inches(0.6), Inches(0.78), size=11, color=WHITE, line_spacing=1.3)

# bottom quote
add_rich(
    s4,
    [
        ("V1 ships ", {"size": 14, "color": RGBColor(0xC0, 0xC0, 0xC0)}),
        ("Mode 1 ", {"size": 14, "color": NEON, "bold": True}),
        ("(direct routing). V2 ships ", {"size": 14, "color": RGBColor(0xC0, 0xC0, 0xC0)}),
        ("pooling + buffer inventory ", {"size": 14, "color": NEON, "bold": True}),
        ("— that's where the moat lives.", {"size": 14, "color": RGBColor(0xC0, 0xC0, 0xC0)}),
    ],
    Inches(0.7), Inches(6.95), Inches(12.2), Inches(0.5),
)

set_speaker_notes(
    s4,
    "Mechanically, iLaunchify is four layers. A Design Studio — a Fabric.js packaging canvas with a live FDA compliance scan, brand asset library, "
    "and print-ready PDF export. A Marketplace — 13 locked product categories with 121 subcategories of manufacturer-built templates filtered by "
    "format, region, capability, and certification. The Orchestration layer — every order decomposes into a multi-partner workflow graph across manufacturer, "
    "label printer, co-packer, warehouse, and the creator's own Shopify or TikTok Shop. And the money-and-manifest layer — Stripe Connect multi-party payouts, "
    "payment held until every partner approves, manifest versioned and audit-logged. V1 ships Mode 1 direct routing. V2 ships pooled production across creators "
    "and a buffer inventory layer — that's where the real moat lives, and that's what scales the unit economics."
)

# =========================================================
# SLIDE 5 — HOW IT WORKS (light)
# =========================================================
s5 = prs.slides.add_slide(BLANK)
add_bg(s5, WHITE)
add_brand_mark(s5, Inches(0.7), Inches(0.5))
add_pill(s5, Inches(0.7), Inches(1.2), Inches(1.4), Inches(0.36), "04 · DEMO", fill=BLACK, text_color=WHITE, size=10)
add_text(
    s5,
    "One creator. One order.\nEverything else hidden.",
    Inches(0.7), Inches(1.75), Inches(12), Inches(1.7),
    font=HEADER_FONT, size=36, bold=True, color=BLACK, line_spacing=1.05,
)

# 3-step flow
step_y = Inches(3.85)
step_h = Inches(2.85)
step_w = Inches(3.95)
gap = Inches(0.2)
sxs = [Inches(0.7), Inches(0.7) + step_w + gap, Inches(0.7) + (step_w + gap) * 2]

steps = [
    ("Pick a template",
     "Creator browses the Marketplace — e.g., a 12-pack kombucha or capsule supplement — pre-validated by a real manufacturer.",
     "[ SCREENSHOT: Marketplace listing ]"),
    ("Customize in Studio",
     "Recipe, label, packaging surfaces. FDA scan runs live as they design. Brand assets pre-fill. Print-ready PDF exports.",
     "[ SCREENSHOT: Design Studio canvas ]"),
    ("We orchestrate",
     "Order auto-routes across partners. Creator sees one timeline. Finished inventory lands in the creator's Shopify or TikTok Shop.",
     "[ SCREENSHOT: Production timeline ]"),
]
for i, (title, body, sshot) in enumerate(steps):
    x = sxs[i]
    add_rect(s5, x, step_y, step_w, step_h, WHITE, line_color=INK_300, line_width=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    # screenshot placeholder area
    add_rect(s5, x + Inches(0.25), step_y + Inches(0.25), step_w - Inches(0.5), Inches(1.15), CREAM, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s5, sshot, x + Inches(0.25), step_y + Inches(0.25), step_w - Inches(0.5), Inches(1.15), size=10, color=INK_500, italic=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # number + title
    add_circle_number(s5, x + Inches(0.3), step_y + Inches(1.55), Inches(0.45), i + 1, fill=PINK, text_color=WHITE)
    add_text(s5, title, x + Inches(0.85), step_y + Inches(1.6), step_w - Inches(1.1), Inches(0.45), size=15, bold=True, color=BLACK, font=HEADER_FONT)
    add_text(s5, body, x + Inches(0.3), step_y + Inches(2.1), step_w - Inches(0.6), Inches(0.75), size=11, color=INK_700, line_spacing=1.3)

# bottom hidden-layer caption
add_rich(
    s5,
    [
        ("The hidden hard part: ", {"size": 12, "bold": True, "color": BLACK, "font": HEADER_FONT}),
        ("payment held until all partners approve. Manifest locked + versioned. Every state transition audit-logged.",
         {"size": 11.5, "color": INK_700}),
    ],
    Inches(0.7), Inches(6.95), Inches(12.2), Inches(0.5),
)

set_speaker_notes(
    s5,
    "The flow is three steps for the creator. Step one — they browse the marketplace and pick a template, say a 12-pack kombucha or a 60-count capsule "
    "supplement. Every template is pre-validated by a real manufacturer with their actual MOQ, lead time, and price ladder. Step two — they customize in "
    "the Design Studio. Recipe, label, packaging surfaces. The FDA compliance scan runs live as they design and tells them in real time when their font is too "
    "small or they're missing an allergen line. Brand assets pre-fill from their library, and they export a print-ready PDF. Step three — we orchestrate. The "
    "order auto-routes across manufacturer, printer, co-packer, and warehouse. The creator sees one unified timeline. The finished inventory lands in their "
    "Shopify or TikTok Shop. Behind that simple flow we hold payment until every partner has approved the manifest, version-lock the artifacts, and write to a "
    "central audit log on every state transition. That's the product. Replace the placeholders with three real screenshots before sending."
)

# =========================================================
# SLIDE 6 — THE MOAT (dark)
# =========================================================
s6 = prs.slides.add_slide(BLANK)
add_bg(s6, BLACK)
add_brand_mark(s6, Inches(0.7), Inches(0.5), on_dark=True)
add_pill(s6, Inches(0.7), Inches(1.2), Inches(1.4), Inches(0.36), "05 · MOAT", fill=NEON, text_color=BLACK, size=10)

add_rich(
    s6,
    [
        ("This is hard to build.\nIt is ", {"size": 38, "color": WHITE, "bold": True, "font": HEADER_FONT}),
        ("harder ", {"size": 38, "color": NEON, "bold": True, "italic": True, "font": "Fraunces"}),
        ("to copy.", {"size": 38, "color": WHITE, "bold": True, "font": HEADER_FONT}),
    ],
    Inches(0.7), Inches(1.75), Inches(12.2), Inches(1.8),
    line_spacing=1.05,
)

# three rows
moats = [
    ("01", "The locked taxonomy",
     "4-layer orthogonal model — 8 creator niches × 13 product categories × manufacturing format × 30 lifestyle tags. The discrimination is taste, not code; you can only build it with deep CPG + creator-economy fluency."),
    ("02", "The 5-layer partner verification",
     "10-state activation FSM. Identity, business, capability, commercial terms, operational standards — all verified before a partner can accept an order. This is operational defensibility. Code is cheap; trusted partner supply isn't."),
    ("03", "The orchestration graph (V2)",
     "Once partner capacity pools across orders, every new creator improves partner economics and every new partner shortens creator lead times. Two-sided network effect, gated by buffer inventory — that's the long-term moat."),
]
mx = Inches(0.7); my = Inches(3.7); mw = Inches(11.95); rh = Inches(1.1); rgap = Inches(0.15)
for i, (n, title, body) in enumerate(moats):
    y = my + (rh + rgap) * i
    add_rect(s6, mx, y, mw, rh, DARK_PANEL, line_color=RGBColor(0x2A, 0x2A, 0x2A), line_width=Pt(0.5), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_rect(s6, mx + Inches(0.3), y + Inches(0.25), Inches(0.6), Inches(0.6), NEON, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s6, n, mx + Inches(0.3), y + Inches(0.25), Inches(0.6), Inches(0.6), size=14, bold=True, color=BLACK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=HEADER_FONT)
    add_text(s6, title, mx + Inches(1.05), y + Inches(0.18), Inches(4), Inches(0.5), size=15, bold=True, color=WHITE, font=HEADER_FONT)
    add_text(s6, body, mx + Inches(1.05), y + Inches(0.55), mw - Inches(1.4), Inches(0.5), size=11, color=RGBColor(0xC8, 0xC8, 0xC8), line_spacing=1.3)

set_speaker_notes(
    s6,
    "Three structural defensibilities. First — the 4-layer marketplace taxonomy. Eight creator niches, 13 product categories, a manufacturing format axis, "
    "30 lifestyle tags. Every layer is orthogonal and intentional. You can copy the database schema, you can't copy the taste-discrimination behind it — that "
    "took deep CPG plus creator-economy fluency to lock. Second — the 5-layer partner verification. A 10-state activation FSM that walks every partner through "
    "identity, business legitimacy, capability, commercial terms, and operational standards verification before they can accept a single order. This is operational "
    "defensibility, not technical. Code is cheap; supply of verified, trustworthy manufacturing partners is the actual scarce resource. Third — the orchestration graph in V2. "
    "Once partner capacity pools across orders and we add a buffer inventory layer, every new creator improves partner economics and every new partner shortens creator "
    "lead times. That's a real two-sided network effect, gated by working capital — the long-term moat."
)

# =========================================================
# SLIDE 7 — MONETIZATION (light)
# =========================================================
s7 = prs.slides.add_slide(BLANK)
add_bg(s7, WHITE)
add_brand_mark(s7, Inches(0.7), Inches(0.5))
add_pill(s7, Inches(0.7), Inches(1.2), Inches(1.7), Inches(0.36), "06 · MONETIZATION", fill=BLACK, text_color=WHITE, size=10)
add_text(
    s7,
    "Three revenue lines.\nLocked, in production.",
    Inches(0.7), Inches(1.75), Inches(12), Inches(1.7),
    font=HEADER_FONT, size=36, bold=True, color=BLACK, line_spacing=1.05,
)

# Tier table
tx = Inches(0.7); ty = Inches(3.7); col_w = Inches(3.95); gap = Inches(0.2); row_h = Inches(2.65)
tiers = [
    ("Maker", "$0/mo", "15%", "Free entry. Unlimited products. Pays platform fee per order."),
    ("Builder", "$49–99/mo*", "12%", "AI design tools. Subscribe & save. Lower per-order fee."),
    ("Agency", "$199–299/mo*", "9%", "Bulk-pricing visibility. Premier partner pool. Sales-touched."),
]
for i, (name, price, fee, body) in enumerate(tiers):
    x = tx + (col_w + gap) * i
    add_rect(s7, x, ty, col_w, row_h, WHITE, line_color=INK_300, line_width=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    # tier-name colored: pink dot beside the name instead of a top accent bar (avoids rounded-corner clipping)
    add_rect(s7, x + Inches(0.3), ty + Inches(0.42), Inches(0.16), Inches(0.16), PINK, shape=MSO_SHAPE.OVAL)
    add_text(s7, name, x + Inches(0.55), ty + Inches(0.25), col_w - Inches(0.85), Inches(0.5), size=22, bold=True, color=BLACK, font=HEADER_FONT)
    add_text(s7, "subscription", x + Inches(0.3), ty + Inches(0.78), col_w - Inches(0.6), Inches(0.3), size=10, color=INK_500)
    add_text(s7, price, x + Inches(0.3), ty + Inches(1.05), col_w - Inches(0.6), Inches(0.45), size=18, bold=True, color=PINK, font=HEADER_FONT)
    add_text(s7, "production fee", x + Inches(0.3), ty + Inches(1.5), col_w - Inches(0.6), Inches(0.3), size=10, color=INK_500)
    add_text(s7, fee, x + Inches(0.3), ty + Inches(1.75), col_w - Inches(0.6), Inches(0.45), size=18, bold=True, color=BLACK, font=HEADER_FONT)
    add_text(s7, body, x + Inches(0.3), ty + Inches(2.2), col_w - Inches(0.6), Inches(0.4), size=10, color=INK_700, line_spacing=1.25)

# revenue lines summary
ry = Inches(6.55)
add_text(s7, "Plus:", Inches(0.7), ry, Inches(0.8), Inches(0.35), size=11, color=INK_500)
add_rich(
    s7,
    [
        ("Marketplace commission ", {"size": 11.5, "bold": True, "color": BLACK}),
        ("on partner payouts (15% / 12% / 8% by partner tier) · ", {"size": 11.5, "color": INK_700}),
        ("Subscribe & Save ", {"size": 11.5, "bold": True, "color": BLACK}),
        ("recurring production for reorders.", {"size": 11.5, "color": INK_700}),
    ],
    Inches(1.4), ry, Inches(11.4), Inches(0.4),
)

# tiny footnote
add_text(s7, "* Builder/Agency monthly price within published bands per PLATFORM_SPEC.md. Annual plans give ~2 months free.",
         Inches(0.7), Inches(7.0), Inches(12.2), Inches(0.3), size=8.5, italic=True, color=INK_500)

set_speaker_notes(
    s7,
    "Three revenue lines, all locked in PLATFORM_SPEC. First — creator subscription. Maker is free with unlimited products and a 15% platform fee per order. "
    "Builder runs $49 to $99 a month with a 12% fee and unlocks AI design and Subscribe and Save. Agency runs $199 to $299 a month with a 9% fee, bulk pricing "
    "visibility, and Premier partner access — sales-touched. Annual plans give roughly two months free. Second — marketplace commission on the partner side: "
    "15, 12, or 8 percent of partner payouts depending on partner tier (Verified, Trusted, Premier). Third — Subscribe and Save: when a consumer reorders on the "
    "creator's channel via our recurring production loop, we capture the recurring production subscription. The schema and Stripe wiring are already built. "
    "Three independent revenue lines means the platform is not single-point-of-failure dependent on subscription, fee, or commission alone."
)

# =========================================================
# SLIDE 8 — BETA PLAN (dark)
# =========================================================
s8 = prs.slides.add_slide(BLANK)
add_bg(s8, BLACK)
add_brand_mark(s8, Inches(0.7), Inches(0.5), on_dark=True)
add_pill(s8, Inches(0.7), Inches(1.2), Inches(1.5), Inches(0.36), "07 · BETA PLAN", fill=NEON, text_color=BLACK, size=10)

add_text(
    s8,
    "We built the platform on our own.\nWe're raising to scale what works.",
    Inches(0.7), Inches(1.75), Inches(12), Inches(1.7),
    font=HEADER_FONT, size=32, bold=True, color=WHITE, line_spacing=1.1,
)

# left column — stat tiles
sx = Inches(0.7); sy = Inches(3.85)
stat_w = Inches(2.65); stat_h = Inches(1.4); s_gap = Inches(0.18)
stats = [
    ("5–8", "creators"),
    ("4–6", "partners"),
    ("90", "days"),
    ("0%", "platform fee"),
]
for i, (big, small) in enumerate(stats):
    x = sx + (stat_w + s_gap) * (i % 2)
    y = sy + (stat_h + s_gap) * (i // 2)
    add_rect(s8, x, y, stat_w, stat_h, DARK_PANEL, line_color=RGBColor(0x2A, 0x2A, 0x2A), line_width=Pt(0.5), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s8, big, x + Inches(0.25), y + Inches(0.18), stat_w - Inches(0.5), Inches(0.7), size=36, bold=True, color=NEON, font=HEADER_FONT)
    add_text(s8, small, x + Inches(0.25), y + Inches(0.9), stat_w - Inches(0.5), Inches(0.4), size=12, color=WHITE)

# right — hard success
rx = Inches(6.7); ry = Inches(3.85); rw = Inches(6.0); rh = Inches(3.0)
add_rect(s8, rx, ry, rw, rh, DARK_PANEL, line_color=RGBColor(0x2A, 0x2A, 0x2A), line_width=Pt(0.5), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s8, "Hard success criteria", rx + Inches(0.3), ry + Inches(0.2), rw - Inches(0.6), Inches(0.4), size=14, bold=True, color=NEON, font=HEADER_FONT)
criteria = [
    "≥4 end-to-end shipments through real channels",
    "≥3 of 8 creators complete the full loop",
    "≥80% partner accept-rate on dispatches",
    "≥75% on-time shipment by committed lead-time",
    "≤2 quality disputes across the cohort",
]
for i, c in enumerate(criteria):
    y = ry + Inches(0.7 + i * 0.42)
    add_rect(s8, rx + Inches(0.3), y + Inches(0.12), Inches(0.12), Inches(0.12), NEON, shape=MSO_SHAPE.OVAL)
    add_text(s8, c, rx + Inches(0.55), y, rw - Inches(0.85), Inches(0.4), size=11.5, color=WHITE)

# bottom narrative band
add_rich(
    s8,
    [
        ("The next pitch becomes: ", {"size": 12.5, "color": RGBColor(0xC0, 0xC0, 0xC0)}),
        ("“We orchestrated $X of production GMV across N partners with on-time-rate Y% and N/M creators reordered.”",
         {"size": 12.5, "color": NEON, "italic": True, "font": "Fraunces"}),
    ],
    Inches(0.7), Inches(6.95), Inches(12.2), Inches(0.45),
)

set_speaker_notes(
    s8,
    "We didn't take outside money to build the platform — we built it on our own. We are running the beta with our own capital. We are raising to scale what works. "
    "The beta is structured: 5 to 8 creators, 4 to 6 partners, 90 days, founder-led white glove, platform fee waived for participants. Hard success criteria are "
    "falsifiable, not vibes — at least 4 end-to-end shipments through real channels, at least 3 of the 8 creators complete the full loop, partner accept-rate at "
    "or above 80%, on-time shipment at or above 75%, and no more than 2 quality disputes. If we hit those, the next fundraise pitch becomes 'we orchestrated this much "
    "real production GMV with this many partners at this on-time rate and this many creators reordered,' which is a categorically different conversation from where "
    "most seed-stage marketplaces start. The de-risk is real."
)

# =========================================================
# SLIDE 9 — MARKET (light)
# =========================================================
s9 = prs.slides.add_slide(BLANK)
add_bg(s9, WHITE)
add_brand_mark(s9, Inches(0.7), Inches(0.5))
add_pill(s9, Inches(0.7), Inches(1.2), Inches(1.4), Inches(0.36), "08 · MARKET", fill=BLACK, text_color=WHITE, size=10)
add_text(
    s9,
    "Big market. Specific wedge.\nDefensible 5-year line.",
    Inches(0.7), Inches(1.75), Inches(12), Inches(1.7),
    font=HEADER_FONT, size=36, bold=True, color=BLACK, line_spacing=1.05,
)

# 3 concentric rectangles (TAM/SAM/SOM)
nx = Inches(0.7); ny = Inches(3.7)
# TAM
add_rect(s9, nx, ny, Inches(6.0), Inches(3.2), CREAM, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s9, "TAM", nx + Inches(0.3), ny + Inches(0.2), Inches(2), Inches(0.4), size=11, bold=True, color=PINK, font=HEADER_FONT)
add_text(s9, "$730B", nx + Inches(0.3), ny + Inches(0.55), Inches(5.4), Inches(0.8), size=42, bold=True, color=BLACK, font=HEADER_FONT)
add_text(s9, "US CPG market — supplements, food & beverage, pet, beauty (Source: FMI / IRI 2024).",
         nx + Inches(0.3), ny + Inches(1.5), Inches(5.4), Inches(0.4), size=10, color=INK_500)
# SAM
add_rect(s9, nx + Inches(0.4), ny + Inches(1.95), Inches(5.2), Inches(1.1), WHITE, line_color=INK_300, line_width=Pt(0.5), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s9, "SAM", nx + Inches(0.6), ny + Inches(2.05), Inches(2), Inches(0.3), size=10, bold=True, color=PINK, font=HEADER_FONT)
add_rich(s9, [
    ("$", {"size": 22, "color": BLACK, "bold": True, "font": HEADER_FONT}),
    ("[VERIFY] B", {"size": 22, "color": BLACK, "bold": True, "font": HEADER_FONT}),
    ("  · creator-led CPG within US (supplements + functional F&B + pet)",
     {"size": 11, "color": INK_500}),
], nx + Inches(0.6), ny + Inches(2.35), Inches(5.0), Inches(0.6))

# Right column — SOM + projection
mx2 = Inches(7.2); my2 = Inches(3.7)
add_rect(s9, mx2, my2, Inches(5.5), Inches(3.2), BLACK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s9, "SOM · 5 YEAR", mx2 + Inches(0.3), my2 + Inches(0.25), Inches(3), Inches(0.4), size=11, bold=True, color=NEON, font=HEADER_FONT)
add_text(s9, "$[VERIFY] M", mx2 + Inches(0.3), my2 + Inches(0.6), Inches(5.0), Inches(0.9), size=44, bold=True, color=WHITE, font=HEADER_FONT)
add_text(
    s9,
    "iLaunchify revenue captured on $[VERIFY] B of creator-led CPG GMV at blended take-rate of ~18% (subscription + production fee + commission).",
    mx2 + Inches(0.3), my2 + Inches(1.6), Inches(5.0), Inches(1.4), size=11, color=RGBColor(0xC0, 0xC0, 0xC0), line_spacing=1.35,
)

# verify note
add_text(
    s9,
    "Numbers marked [VERIFY] need defensible public sourcing (IRI / Statista / NielsenIQ / SignalFire creator-economy report) before sending.",
    Inches(0.7), Inches(7.0), Inches(12.2), Inches(0.4),
    size=9, italic=True, color=INK_500,
)

set_speaker_notes(
    s9,
    "Market sizing — honest. The US CPG market is roughly 730 billion dollars (FMI and IRI 2024 baseline). That's the TAM. Our serviceable available market is "
    "the creator-led slice of US CPG within the V1 categories — supplements, functional food and beverage, and pet — which is a much smaller and faster-growing "
    "sub-segment. The 5-year SOM is the production-orchestration revenue we capture on the creator-led CPG GMV that runs through iLaunchify, at a blended take "
    "rate that combines subscription, platform fee, and marketplace commission — historically around 18 percent. Pavel: I left the SAM and SOM numbers as placeholders "
    "because I want defensible public sources behind them — SignalFire's creator economy report, Statista, NielsenIQ — rather than something I made up. Fill those in "
    "from real public sources before sending. The TAM line is solid."
)

# =========================================================
# SLIDE 10 — PRODUCT VELOCITY (light)
# =========================================================
s10 = prs.slides.add_slide(BLANK)
add_bg(s10, WHITE)
add_brand_mark(s10, Inches(0.7), Inches(0.5))
add_pill(s10, Inches(0.7), Inches(1.2), Inches(2.05), Inches(0.36), "09 · WHY US, WHY NOW", fill=BLACK, text_color=WHITE, size=10)

add_rich(
    s10,
    [
        ("Six months. ", {"size": 36, "color": BLACK, "bold": True, "font": HEADER_FONT}),
        ("One founder.\nA real platform.", {"size": 36, "color": PINK, "bold": True, "font": HEADER_FONT, "italic": True}),
    ],
    Inches(0.7), Inches(1.75), Inches(12), Inches(1.8),
    line_spacing=1.05,
)

# big stat row
sy = Inches(3.75)
sw = Inches(2.95); sh_ = Inches(1.3); sg = Inches(0.15)
sxs = [Inches(0.7), Inches(0.7) + sw + sg, Inches(0.7) + (sw + sg) * 2, Inches(0.7) + (sw + sg) * 3]
big_stats = [
    ("~590", "shipped tasks"),
    ("4", "Next.js apps"),
    ("100+", "Prisma models"),
    ("24", "DB migrations"),
]
for i, (big, small) in enumerate(big_stats):
    x = sxs[i]
    add_rect(s10, x, sy, sw, sh_, CREAM, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s10, big, x + Inches(0.3), sy + Inches(0.15), sw - Inches(0.6), Inches(0.7), size=34, bold=True, color=BLACK, font=HEADER_FONT)
    add_text(s10, small, x + Inches(0.3), sy + Inches(0.85), sw - Inches(0.6), Inches(0.4), size=11, color=INK_500)

# bullet list of what shipped (2 columns)
bx = Inches(0.7); by = Inches(5.3); bw = Inches(5.95); bgap = Inches(0.15)
left_items = [
    "Fabric.js Design Studio with live FDA compliance scan",
    "Stripe Connect multi-party payouts (creator + partner)",
    "5-layer partner onboarding + 10-state activation FSM",
]
right_items = [
    "Multi-partner approval workflow with manifest versioning",
    "Marketplace 4-layer taxonomy + niche auto-suggest engine",
    "Lawyer-ready legal foundation (ToS / Privacy / 2 agreements)",
]

for col, items in enumerate([left_items, right_items]):
    cx = bx + (bw + bgap) * col
    for i, item in enumerate(items):
        ry = by + Inches(i * 0.45)
        add_rect(s10, cx, ry + Inches(0.13), Inches(0.12), Inches(0.12), PINK, shape=MSO_SHAPE.OVAL)
        add_text(s10, item, cx + Inches(0.25), ry, bw - Inches(0.3), Inches(0.4), size=11.5, color=INK_700)

# closer
add_rich(
    s10,
    [
        ("This is what one focused founder built. ", {"size": 13, "color": INK_500, "italic": True, "font": "Fraunces"}),
        ("Imagine what a team does.", {"size": 13, "color": BLACK, "bold": True, "font": HEADER_FONT}),
    ],
    Inches(0.7), Inches(6.95), Inches(12), Inches(0.4),
)

set_speaker_notes(
    s10,
    "The proof point that matters to a seed investor — velocity. In roughly six months of focused work, one founder shipped four Next.js applications, around 590 "
    "tracked tasks, 100+ Prisma models on a CockroachDB Serverless backend, and 24 production migrations. Inside that we shipped a Fabric.js packaging Design Studio "
    "with a live FDA compliance scan, full Stripe Connect multi-party payouts on both sides, a 5-layer partner onboarding with a 10-state activation FSM, a multi-partner "
    "approval workflow with manifest versioning, the marketplace taxonomy with deterministic niche auto-suggestion, and lawyer-ready legal documents. The architecture is "
    "locked, the gotchas are gotcha'd, and the team-of-one velocity is what compounds when we hire. The next dollar gets a lot of leverage."
)

# =========================================================
# SLIDE 11 — TEAM (light)
# =========================================================
s11 = prs.slides.add_slide(BLANK)
add_bg(s11, WHITE)
add_brand_mark(s11, Inches(0.7), Inches(0.5))
add_pill(s11, Inches(0.7), Inches(1.2), Inches(1.3), Inches(0.36), "10 · TEAM", fill=BLACK, text_color=WHITE, size=10)
add_text(
    s11,
    "Solo today.\nCompounding by month 6.",
    Inches(0.7), Inches(1.75), Inches(12), Inches(1.7),
    font=HEADER_FONT, size=36, bold=True, color=BLACK, line_spacing=1.05,
)

# Founder card
fx = Inches(0.7); fy = Inches(3.7); fw = Inches(5.8); fh = Inches(3.2)
add_rect(s11, fx, fy, fw, fh, CREAM, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
# avatar placeholder circle
add_rect(s11, fx + Inches(0.4), fy + Inches(0.4), Inches(1.3), Inches(1.3), PINK, shape=MSO_SHAPE.OVAL)
add_text(s11, "PG", fx + Inches(0.4), fy + Inches(0.4), Inches(1.3), Inches(1.3), size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=HEADER_FONT)
add_text(s11, "Pavel Georgiev", fx + Inches(1.9), fy + Inches(0.5), Inches(3.5), Inches(0.5), size=20, bold=True, color=BLACK, font=HEADER_FONT)
add_text(s11, "Founder & CEO", fx + Inches(1.9), fy + Inches(1.0), Inches(3.5), Inches(0.4), size=12, color=PINK, bold=True)
add_text(
    s11,
    "[BIO PLACEHOLDER — prior CPG-adjacent operating role at FOD; deep familiarity with the contract manufacturing landscape; full-stack founder shipping every layer of the V1 platform himself.]",
    fx + Inches(0.4), fy + Inches(1.95), fw - Inches(0.8), Inches(1.15), size=11, italic=True, color=INK_700, line_spacing=1.35,
)

# Right column: stack + hiring plan
rx = Inches(6.85); ry = Inches(3.7); rw = Inches(5.85); rh = Inches(3.2)
add_rect(s11, rx, ry, rw, rh, WHITE, line_color=INK_300, line_width=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s11, "STACK", rx + Inches(0.3), ry + Inches(0.25), Inches(2), Inches(0.3), size=10, bold=True, color=PINK, font=HEADER_FONT)
add_text(s11, "Next.js 15 · React 19 · TypeScript strict · CockroachDB Serverless · Prisma · Stripe Connect · Cloudflare R2 · Fabric.js · Auth.js v5 · Resend.",
         rx + Inches(0.3), ry + Inches(0.55), rw - Inches(0.6), Inches(1.1), size=11.5, color=INK_700, line_spacing=1.35)

add_text(s11, "HIRING WITH THIS ROUND", rx + Inches(0.3), ry + Inches(1.75), Inches(4), Inches(0.3), size=10, bold=True, color=PINK, font=HEADER_FONT)
hires = [
    "Senior product engineer (full-stack)",
    "Partner ops manager (CPG manufacturing)",
    "Product / brand designer",
]
for i, h in enumerate(hires):
    hy = ry + Inches(2.1 + i * 0.32)
    add_rect(s11, rx + Inches(0.3), hy + Inches(0.1), Inches(0.12), Inches(0.12), BLACK, shape=MSO_SHAPE.OVAL)
    add_text(s11, h, rx + Inches(0.55), hy, rw - Inches(0.85), Inches(0.35), size=11.5, color=INK_700)

# bottom italic
add_rich(
    s11,
    [
        ("Architecture is locked. ", {"size": 13, "color": INK_500, "italic": True, "font": "Fraunces"}),
        ("The team won't redesign the foundation — they'll compound on it.", {"size": 13, "color": BLACK, "bold": True}),
    ],
    Inches(0.7), Inches(6.95), Inches(12), Inches(0.4),
)

set_speaker_notes(
    s11,
    "Honesty up front — I am the only person on the team today. That's why the architectural moat slide matters: the foundation is locked correctly, so the people I hire "
    "with this round are not going to spend their first 6 months refactoring my code; they're going to compound on it. The stack is mainstream and boring on purpose: "
    "Next.js 15, React 19, strict TypeScript, CockroachDB Serverless via Prisma, Stripe Connect, Cloudflare R2, Fabric.js, Auth.js v5, Resend. Nothing exotic, nothing "
    "that requires hiring rare specialists. Hiring plan in the first six months — one senior product engineer to own velocity on the surface area I built, one partner "
    "operations manager with CPG manufacturing relationships and the skin-deep familiarity to onboard partners faster than I can, and one product or brand designer. "
    "Pavel: replace the [BIO PLACEHOLDER] line with your actual short bio before sending."
)

# =========================================================
# SLIDE 12 — THE ASK (dark)
# =========================================================
s12 = prs.slides.add_slide(BLANK)
add_bg(s12, BLACK)
add_brand_mark(s12, Inches(0.7), Inches(0.5), on_dark=True)
add_pill(s12, Inches(0.7), Inches(1.2), Inches(1.4), Inches(0.36), "11 · THE ASK", fill=NEON, text_color=BLACK, size=10)

add_rich(
    s12,
    [
        ("We're raising ", {"size": 40, "color": WHITE, "bold": True, "font": HEADER_FONT}),
        ("[$X.X M]", {"size": 40, "color": NEON, "bold": True, "font": HEADER_FONT}),
        ("\nseed.", {"size": 40, "color": WHITE, "bold": True, "font": HEADER_FONT}),
    ],
    Inches(0.7), Inches(1.75), Inches(12), Inches(1.8),
    line_spacing=1.05,
)

# Use of funds — pie-style horizontal bars
uy = Inches(3.85)
use_of_funds = [
    ("Engineering · 2 hires + infra hardening", 50, NEON),
    ("Partner network expansion (4 service types, top 3 markets)", 25, PINK),
    ("Creator acquisition + brand", 15, WHITE),
    ("Legal · compliance counsel · admin", 10, RGBColor(0x6B, 0x6B, 0x6B)),
]
bar_x = Inches(0.7); bar_w_full = Inches(6.5); bar_h = Inches(0.4); bar_gap = Inches(0.55)
add_text(s12, "Use of funds", bar_x, uy - Inches(0.45), Inches(4), Inches(0.4), size=12, bold=True, color=NEON, font=HEADER_FONT)
for i, (label, pct, color) in enumerate(use_of_funds):
    y = uy + (bar_gap) * i
    # full track
    add_rect(s12, bar_x, y + Inches(0.18), bar_w_full, Inches(0.18), DARK_PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    # filled portion
    add_rect(s12, bar_x, y + Inches(0.18), Emu(int(bar_w_full * pct / 100)), Inches(0.18), color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s12, f"{pct}%", bar_x + bar_w_full + Inches(0.15), y + Inches(0.1), Inches(0.7), Inches(0.35), size=12, bold=True, color=WHITE, font=HEADER_FONT)
    add_text(s12, label, bar_x, y - Inches(0.05), bar_w_full, Inches(0.3), size=11, color=RGBColor(0xD0, 0xD0, 0xD0))

# Milestones panel on right
my_ = Inches(3.4); mx_ = Inches(8.1); mw_ = Inches(4.6); mh_ = Inches(3.6)
add_rect(s12, mx_, my_, mw_, mh_, DARK_PANEL, line_color=RGBColor(0x2A, 0x2A, 0x2A), line_width=Pt(0.5), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s12, "18-MONTH MILESTONES", mx_ + Inches(0.3), my_ + Inches(0.25), mw_ - Inches(0.6), Inches(0.4), size=11, bold=True, color=NEON, font=HEADER_FONT)

milestones = [
    ("M3", "Beta closes with documented transaction data + first case studies"),
    ("M6", "GA · 50 active creators · 20 active partners"),
    ("M12", "$[VERIFY] M annualized GMV across the platform"),
    ("M18", "Mode 1 → Mode 2 transition (pooled production live)"),
]
for i, (when, what) in enumerate(milestones):
    y = my_ + Inches(0.85 + i * 0.65)
    # tag pill
    add_rect(s12, mx_ + Inches(0.3), y, Inches(0.65), Inches(0.4), NEON, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s12, when, mx_ + Inches(0.3), y, Inches(0.65), Inches(0.4), size=11, bold=True, color=BLACK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=HEADER_FONT)
    add_text(s12, what, mx_ + Inches(1.05), y - Inches(0.02), mw_ - Inches(1.35), Inches(0.55), size=11, color=WHITE, line_spacing=1.3)

# bottom CTA band
add_rich(
    s12,
    [
        ("Let's talk. ", {"size": 14, "color": NEON, "bold": True, "italic": True, "font": "Fraunces"}),
        ("georgiev.pavel@gmail.com  ·  ilaunchify.com", {"size": 14, "color": WHITE}),
    ],
    Inches(0.7), Inches(7.0), Inches(12), Inches(0.4),
)

set_speaker_notes(
    s12,
    "The ask. We are raising a seed round of [Pavel — set the actual number; suggested $2.0M as a reasonable mid-anchor for this stage, typical band is $1M–$3M]. "
    "Allocation: 50% engineering — two senior hires plus production infrastructure hardening; 25% partner network expansion across the four service types in the top "
    "three US markets, founder-led on the ground; 15% creator acquisition and brand work; 10% legal, compliance counsel, and admin overhead. Milestones over 18 months: "
    "by month 3 the beta closes with documented transaction data and first case studies; by month 6 we are at GA with 50 active creators and 20 active partners; by "
    "month 12 we are at $[VERIFY]M annualized GMV across the platform; by month 18 we begin the transition from Mode 1 direct routing to Mode 2 pooled production — "
    "the moat features go live. Set the dollar amount and the M12 GMV target to your actual numbers. The conversation we want is the follow-on diligence call. My email "
    "is on the slide."
)

# write
import os
# When running inside the workspace VM, map to the mounted path. Local run uses the host path.
candidates = [
    "/Users/soundstation/Documents/CLAUDE/iLaunchify/docs/pitch",
    "/sessions/brave-affectionate-mayer/mnt/iLaunchify/docs/pitch",
]
out_dir = next((c for c in candidates if os.path.isdir(os.path.dirname(c)) or os.path.isdir(c)), candidates[-1])
out_path = os.path.join(out_dir, "iLaunchify_Seed_Deck.pptx")
os.makedirs(out_dir, exist_ok=True)
prs.save(out_path)
print(f"WROTE {out_path}")
print(f"Slide count: {len(prs.slides)}")
